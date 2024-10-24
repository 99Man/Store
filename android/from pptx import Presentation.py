from pptx import Presentation

# Create a PowerPoint Presentation object
presentation = Presentation()


# Function to add a slide with a title and content
def add_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[
        1
    ]  # Use the slide layout with title and content
    slide = presentation.slides.add_slide(slide_layout)

    # Add title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add content
    content_placeholder = slide.shapes.placeholders[1]
    content_placeholder.text = content


# Add slides
add_slide(
    presentation,
    "Is Social Media Harmful for Teenagers?",
    "Overview: Exploring arguments for and against the harmful effects of social media on teenagers.",
)

add_slide(
    presentation,
    "Arguments for Social Media Being Harmful",
    "1. Mental Health Issues: Anxiety, depression, cyberbullying.\n"
    "2. Impact on Self-Esteem: Unrealistic beauty standards, validation.\n"
    "3. Addiction and Distraction: Decreased productivity, academic impact.",
)

add_slide(
    presentation,
    "Counterarguments: Benefits of Social Media",
    "1. Connection and Community: Builds relationships, support networks.\n"
    "2. Access to Information: Mobilizing movements, social issues awareness.\n"
    "3. Creativity and Self-Expression: Showcasing talents, content creation.",
)

add_slide(
    presentation,
    "Conclusion and Call to Action",
    "Summary: Social media's impact on teenagers is complex.\n"
    "Call to Action: Encourage discussions and educate on responsible use.",
)

# Save the presentation
presentation.save("Social_Media_Harmful_for_Teenagers.pptx")
